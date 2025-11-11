# imap_email_ingestion_pipeline/attachment_processor.py
# Multi-format attachment processor with intelligent OCR fallback
# Handles PDFs, Office documents, images with confidence scoring
# RELEVANT FILES: ocr_engine.py, state_manager.py, entity_extractor.py

import os
import tempfile
import logging
from pathlib import Path
from typing import Dict, List, Any, Optional, Tuple
import hashlib
import magic
import json
from datetime import datetime

# Document processing imports
try:
    import PyPDF2
    import openpyxl
    from docx import Document
    from pptx import Presentation
    import tabula
except ImportError as e:
    logging.warning(f"Some document processing libraries not available: {e}")

# OCR engine import
try:
    from ocr_engine import OCREngine
except ImportError:
    OCREngine = None

class AttachmentProcessor:
    def __init__(self, storage_path: str = "./data/attachments"):
        self.storage_path = Path(storage_path)
        self.storage_path.mkdir(parents=True, exist_ok=True)
        self.logger = logging.getLogger(__name__)
        
        # Initialize OCR engine if available
        self.ocr_engine = OCREngine() if OCREngine else None

        # Check Tabula availability for table extraction
        self.tabula_available = self._check_tabula_available()
        if self.tabula_available:
            self.logger.info("Tabula table extraction enabled")

        # File size limits (in bytes)
        self.max_file_size = 100 * 1024 * 1024  # 100MB
        self.max_pdf_pages = 50

        # Resource limits to prevent memory exhaustion
        self.max_excel_rows = 5000  # Limit Excel rows to prevent OOM
        self.max_excel_sheets = 20  # Limit number of sheets to process
        self.max_word_paragraphs = 10000  # Limit Word doc paragraphs
        self.max_ppt_slides = 100  # Limit PowerPoint slides
        
        # Supported file types
        self.supported_types = {
            'application/pdf': self._process_pdf,
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': self._process_excel,
            'application/vnd.ms-excel': self._process_excel,
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document': self._process_word,
            'application/msword': self._process_word,
            'application/vnd.openxmlformats-officedocument.presentationml.presentation': self._process_powerpoint,
            'application/vnd.ms-powerpoint': self._process_powerpoint,
            'image/jpeg': self._process_image,
            'image/png': self._process_image,
            'image/tiff': self._process_image,
            'image/bmp': self._process_image,
            'text/plain': self._process_text,
            'text/csv': self._process_csv
        }

    def _check_tabula_available(self) -> bool:
        """Check if tabula-py is available for table extraction"""
        try:
            import tabula
            return True
        except ImportError:
            self.logger.info("Tabula not available - table extraction disabled")
            return False

    def process_attachment(self, attachment_data: Dict[str, Any], email_uid: str) -> Dict[str, Any]:
        """Process single attachment and return extracted content"""
        try:
            # Extract attachment content
            content = attachment_data['part'].get_payload(decode=True)
            if not content:
                return {'error': 'No content in attachment'}
            
            # Check file size
            if len(content) > self.max_file_size:
                return {
                    'error': f'File too large: {len(content)} bytes > {self.max_file_size}',
                    'filename': attachment_data['filename'],
                    'size': len(content)
                }
            
            # Generate file hash for deduplication
            file_hash = hashlib.sha256(content).hexdigest()
            
            # Detect actual MIME type
            try:
                mime_type = magic.from_buffer(content, mime=True)
            except (AttributeError, TypeError, Exception) as e:
                # Magic library can raise various exceptions, fallback to content_type from email
                self.logger.debug(f"Magic MIME detection failed: {e}, using email content_type")
                mime_type = attachment_data.get('content_type', 'application/octet-stream')
            
            # Create storage directory
            storage_dir = self._create_storage_directory(email_uid, file_hash)

            # Check if extraction already exists (caching to avoid redundant processing)
            extracted_txt_path = storage_dir / 'extracted.txt'
            metadata_path = storage_dir / 'metadata.json'

            if extracted_txt_path.exists() and metadata_path.exists():
                # Load cached result
                try:
                    with open(extracted_txt_path, 'r', encoding='utf-8') as f:
                        cached_text = f.read()

                    with open(metadata_path, 'r', encoding='utf-8') as f:
                        cached_metadata = json.load(f)

                    self.logger.info(f"Using cached extraction for {attachment_data['filename']} (hash: {file_hash[:16]}...)")

                    return {
                        'filename': attachment_data['filename'],
                        'file_hash': file_hash,
                        'mime_type': cached_metadata.get('file_info', {}).get('mime_type', mime_type),
                        'file_size': len(content),
                        'storage_path': str(storage_dir),
                        'processing_status': 'completed',
                        'processing_method': 'cached',
                        'extracted_text': cached_text,
                        'data': cached_metadata.get('processing', {}).get('tables', {}),
                        'cached': True
                    }
                except Exception as e:
                    self.logger.warning(f"Failed to load cached extraction: {e}, reprocessing...")
                    # Continue with normal processing if cache load fails

            # Save original file
            original_path = storage_dir / 'original' / attachment_data['filename']
            original_path.parent.mkdir(parents=True, exist_ok=True)

            with open(original_path, 'wb') as f:
                f.write(content)

            # Process based on file type
            processing_result = self._process_by_type(
                content, mime_type, attachment_data['filename'], storage_dir
            )
            
            # Build result dictionary
            result = {
                'filename': attachment_data['filename'],
                'file_hash': file_hash,
                'mime_type': mime_type,
                'file_size': len(content),
                'storage_path': str(storage_dir),
                'processing_status': processing_result.get('status', 'completed'),
                'extraction_method': processing_result.get('method', 'unknown'),
                'extracted_text': processing_result.get('text', ''),
                'extracted_data': processing_result.get('data', {}),
                'ocr_confidence': processing_result.get('ocr_confidence', 0.0),
                'page_count': processing_result.get('page_count', 1),
                'error': processing_result.get('error', None)
            }
            
            # Save extracted content
            if result['extracted_text']:
                extracted_path = storage_dir / 'extracted.txt'
                with open(extracted_path, 'w', encoding='utf-8') as f:
                    f.write(result['extracted_text'])

            # Create metadata.json for traceability
            try:
                metadata = {
                    "source_type": "email_attachment",
                    "source_context": {
                        "email_uid": email_uid,
                        "email_subject": attachment_data.get('email_subject', 'Unknown'),
                        "email_date": attachment_data.get('email_date', datetime.now().isoformat())
                    },
                    "file_info": {
                        "original_filename": attachment_data['filename'],
                        "file_hash": file_hash,
                        "file_size": len(content),
                        "mime_type": mime_type
                    },
                    "processing": {
                        "timestamp": datetime.now().isoformat(),
                        "extraction_method": processing_result.get('method', 'unknown'),
                        "status": processing_result.get('status', 'completed'),
                        "page_count": processing_result.get('page_count', 1),
                        "text_chars": len(processing_result.get('text', '')),
                        "tables_extracted": len(processing_result.get('data', {}).get('tables', [])),
                        "ocr_confidence": processing_result.get('ocr_confidence', 0.0)
                    },
                    "storage": {
                        "original_path": f"original/{attachment_data['filename']}",
                        "extracted_text_path": "extracted.txt" if result['extracted_text'] else None,
                        "created_at": datetime.now().isoformat()
                    }
                }

                metadata_path = storage_dir / 'metadata.json'
                with open(metadata_path, 'w', encoding='utf-8') as f:
                    json.dump(metadata, f, indent=2)

                self.logger.debug(f"Saved metadata for {attachment_data['filename']}")

            except Exception as e:
                # Log but don't fail the entire process
                self.logger.warning(f"Failed to create metadata.json for {attachment_data['filename']}: {e}")

            self.logger.info(f"Processed {attachment_data['filename']}: "
                           f"{len(result['extracted_text'])} chars extracted")
            
            return result
            
        except Exception as e:
            self.logger.error(f"Error processing attachment {attachment_data.get('filename', 'unknown')}: {e}")
            return {
                'error': str(e),
                'filename': attachment_data.get('filename', 'unknown'),
                'processing_status': 'failed'
            }
    
    def _create_storage_directory(self, email_uid: str, file_hash: str) -> Path:
        """Create storage directory structure using email_uid for isolation"""
        # FIX: Use email_uid to prevent file collision when same attachment in multiple emails
        # Pattern matches DoclingProcessor implementation for consistency
        dir_path = self.storage_path / email_uid / file_hash
        dir_path.mkdir(parents=True, exist_ok=True)
        return dir_path
    
    def _process_by_type(self, content: bytes, mime_type: str, filename: str,
                        storage_dir: Path) -> Dict[str, Any]:
        """Route processing based on file type"""
        try:
            if mime_type in self.supported_types:
                processor = self.supported_types[mime_type]
                return processor(content, filename, storage_dir)
            else:
                # Unknown type - try OCR if it's an image-like file
                if self.ocr_engine and any(ext in filename.lower() for ext in ['.jpg', '.jpeg', '.png', '.tiff', '.bmp']):
                    self.logger.info(f"📸 Processing image file with OCR: {filename}")
                    return self._process_image(content, filename, storage_dir)
                else:
                    # Log explicitly which file types are supported
                    file_ext = Path(filename).suffix.lower()
                    self.logger.warning(
                        f"⚠️ Unsupported file type for {filename} (mime: {mime_type}, ext: {file_ext})\n"
                        f"   Supported types: PDF, Excel (xlsx/xls), Word (docx/doc), PowerPoint (pptx/ppt), Images (with OCR)"
                    )
                    return {
                        'status': 'unsupported',
                        'method': 'none',
                        'text': '',
                        'error': f'Unsupported file type: {mime_type} ({file_ext}). Supported: PDF, Excel, Word, PowerPoint, Images'
                    }
        except Exception as e:
            self.logger.error(f"❌ Failed to process {filename}: {e}")
            return {
                'status': 'error',
                'method': 'failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_pdf(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process PDF files with OCR fallback"""
        try:
            # Save content to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pdf', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                # Try native PDF text extraction first
                extracted_text = ""
                page_count = 0
                
                with open(tmp_path, 'rb') as pdf_file:
                    pdf_reader = PyPDF2.PdfReader(pdf_file)
                    page_count = len(pdf_reader.pages)
                    
                    # Limit pages to prevent excessive processing
                    max_pages = min(page_count, self.max_pdf_pages)
                    
                    for page_num in range(max_pages):
                        try:
                            page = pdf_reader.pages[page_num]
                            page_text = page.extract_text()
                            if page_text:
                                extracted_text += page_text + "\n"
                        except Exception as e:
                            self.logger.warning(f"Failed to extract text from page {page_num}: {e}")

                # Extract tables with Tabula (graceful: returns [] if unavailable/fails)
                tables = self._extract_tables_tabula(tmp_path)

                # Check if we got meaningful text
                if len(extracted_text.strip()) > 50:
                    # Native extraction successful
                    return {
                        'status': 'completed',
                        'method': 'pdf_native',
                        'text': extracted_text.strip(),
                        'data': {'tables': tables},
                        'page_count': page_count,
                        'ocr_confidence': 1.0
                    }
                else:
                    # Try OCR fallback
                    if self.ocr_engine:
                        self.logger.info(f"PDF {filename} has minimal text, trying OCR...")
                        ocr_result = self.ocr_engine.process_pdf(tmp_path)
                        return {
                            'status': 'completed',
                            'method': 'pdf_ocr',
                            'text': ocr_result.get('text', ''),
                            'data': {'tables': []},
                            'page_count': page_count,
                            'ocr_confidence': ocr_result.get('confidence', 0.0)
                        }
                    else:
                        return {
                            'status': 'partial',
                            'method': 'pdf_native',
                            'text': extracted_text.strip(),
                            'data': {'tables': []},
                            'page_count': page_count,
                            'ocr_confidence': 0.5,
                            'warning': 'Minimal text extracted, OCR not available'
                        }
                
            finally:
                # Clean up temp file
                os.unlink(tmp_path)
                
        except Exception as e:
            return {
                'status': 'error',
                'method': 'pdf_failed',
                'text': '',
                'error': str(e)
            }

    def _extract_tables_tabula(self, pdf_path: str) -> List[Dict[str, Any]]:
        """
        Extract tables using Tabula (Java-based table extraction)

        Graceful degradation: Returns empty list if Tabula unavailable or fails.
        Format matches Docling's table structure for consistency.
        """
        if not self.tabula_available:
            return []

        try:
            import tabula
            # FIX: Limit table extraction to same pages as text extraction for consistency
            dfs = tabula.read_pdf(pdf_path, pages=f'1-{self.max_pdf_pages}', multiple_tables=True, silent=True)

            tables = []
            for idx, df in enumerate(dfs):
                tables.append({
                    'index': idx,
                    'data': df.to_dict(orient='records'),
                    'num_rows': len(df),
                    'num_cols': len(df.columns),
                    'error': None
                })

            if tables:
                self.logger.info(f"Tabula extracted {len(tables)} table(s)")
            else:
                self.logger.info("Tabula found no tables in PDF")

            return tables

        except Exception as e:
            self.logger.warning(f"Tabula extraction failed: {e}")
            return []

    def _extract_tables_excel(self, workbook) -> List[Dict[str, Any]]:
        """
        Extract structured tables from Excel workbook.

        FIX: Returns tables in TableEntityExtractor-compatible format.
        Pattern matches _extract_tables_tabula() for consistency with PDF processing.
        """
        import pandas as pd
        tables = []
        table_idx = 0

        try:
            # Limit number of sheets to prevent memory exhaustion
            sheets_to_process = workbook.sheetnames[:self.max_excel_sheets]
            if len(workbook.sheetnames) > self.max_excel_sheets:
                self.logger.warning(f"Excel file has {len(workbook.sheetnames)} sheets, processing only first {self.max_excel_sheets}")

            for sheet_name in sheets_to_process:
                worksheet = workbook[sheet_name]

                # Convert worksheet to list of lists (rows) with row limit
                rows = []
                row_count = 0
                for row in worksheet.iter_rows(values_only=True):
                    # Filter out completely empty rows
                    if any(cell is not None for cell in row):
                        rows.append([str(cell) if cell is not None else '' for cell in row])
                        row_count += 1
                        # Enforce row limit to prevent memory exhaustion
                        if row_count >= self.max_excel_rows:
                            self.logger.warning(f"Sheet '{sheet_name}' has more than {self.max_excel_rows} rows, truncating")
                            break

                if len(rows) < 2:  # Need at least header + 1 data row
                    continue

                # Create DataFrame with first row as header
                df = pd.DataFrame(rows[1:], columns=rows[0])

                # Filter out empty columns
                df = df.loc[:, (df != '').any(axis=0)]

                if df.empty or len(df.columns) == 0:
                    continue

                tables.append({
                    'index': table_idx,
                    'data': df.to_dict(orient='records'),
                    'num_rows': len(df),
                    'num_cols': len(df.columns),
                    'sheet_name': sheet_name,
                    'error': None
                })
                table_idx += 1

            if tables:
                self.logger.info(f"Extracted {len(tables)} table(s) from Excel")
            return tables

        except Exception as e:
            self.logger.warning(f"Excel table extraction failed: {e}")
            return []

    def _extract_tables_word(self, doc) -> List[Dict[str, Any]]:
        """
        Extract structured tables from Word document.

        FIX: Returns tables in TableEntityExtractor-compatible format.
        Pattern matches _extract_tables_tabula() for consistency with PDF processing.
        """
        import pandas as pd
        tables = []

        try:
            for idx, table in enumerate(doc.tables):
                # Extract table as list of lists
                rows = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    # Filter out completely empty rows
                    if any(cell for cell in row_data):
                        rows.append(row_data)

                if len(rows) < 2:  # Need at least header + 1 data row
                    continue

                # Create DataFrame with first row as header
                df = pd.DataFrame(rows[1:], columns=rows[0])

                # Filter out empty columns
                df = df.loc[:, (df != '').any(axis=0)]

                if df.empty or len(df.columns) == 0:
                    continue

                tables.append({
                    'index': idx,
                    'data': df.to_dict(orient='records'),
                    'num_rows': len(df),
                    'num_cols': len(df.columns),
                    'error': None
                })

            if tables:
                self.logger.info(f"Extracted {len(tables)} table(s) from Word")
            return tables

        except Exception as e:
            self.logger.warning(f"Word table extraction failed: {e}")
            return []

    def _extract_tables_powerpoint(self, prs) -> List[Dict[str, Any]]:
        """
        Extract structured tables from PowerPoint presentation.

        FIX: Returns tables in TableEntityExtractor-compatible format.
        Pattern matches _extract_tables_tabula() for consistency with PDF processing.
        """
        import pandas as pd
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        tables = []
        table_idx = 0

        try:
            for slide_num, slide in enumerate(prs.slides):
                for shape in slide.shapes:
                    # Check if shape is a table
                    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                        # Extract table as list of lists
                        rows = []
                        for row in shape.table.rows:
                            row_data = [cell.text.strip() for cell in row.cells]
                            # Filter out completely empty rows
                            if any(cell for cell in row_data):
                                rows.append(row_data)

                        if len(rows) < 2:  # Need at least header + 1 data row
                            continue

                        # Create DataFrame with first row as header
                        df = pd.DataFrame(rows[1:], columns=rows[0])

                        # Filter out empty columns
                        df = df.loc[:, (df != '').any(axis=0)]

                        if df.empty or len(df.columns) == 0:
                            continue

                        tables.append({
                            'index': table_idx,
                            'data': df.to_dict(orient='records'),
                            'num_rows': len(df),
                            'num_cols': len(df.columns),
                            'slide_number': slide_num + 1,
                            'error': None
                        })
                        table_idx += 1

            if tables:
                self.logger.info(f"Extracted {len(tables)} table(s) from PowerPoint")
            return tables

        except Exception as e:
            self.logger.warning(f"PowerPoint table extraction failed: {e}")
            return []

    def _process_excel(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process Excel files extracting data and formulas"""
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                # FIX: Load with data_only=True to get computed values instead of formulas
                workbook = openpyxl.load_workbook(tmp_path, data_only=True)
                extracted_text = []

                # Extract text for content search with limits
                sheets_to_process = workbook.sheetnames[:self.max_excel_sheets]
                for sheet_name in sheets_to_process:
                    worksheet = workbook[sheet_name]

                    row_count = 0
                    for row in worksheet.iter_rows():
                        for cell in row:
                            if cell.value is not None:
                                cell_info = f"{cell.coordinate}: {cell.value}"
                                extracted_text.append(cell_info)
                        row_count += 1
                        # Enforce row limit for text extraction
                        if row_count >= self.max_excel_rows:
                            break

                # FIX: Extract structured tables using helper (TableEntityExtractor-compatible format)
                tables = self._extract_tables_excel(workbook)

                return {
                    'status': 'completed',
                    'method': 'excel_native',
                    'text': '\n'.join(extracted_text),
                    'data': {'tables': tables},  # FIX: Unified format (was 'worksheets')
                    'ocr_confidence': 1.0
                }
                
            finally:
                os.unlink(tmp_path)
                
        except Exception as e:
            return {
                'status': 'error',
                'method': 'excel_failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_word(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process Word documents"""
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.docx', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                doc = Document(tmp_path)
                extracted_text = []

                # Extract paragraphs with limit to prevent memory exhaustion
                paragraph_count = 0
                for paragraph in doc.paragraphs:
                    if paragraph.text.strip():
                        extracted_text.append(paragraph.text)
                        paragraph_count += 1
                        if paragraph_count >= self.max_word_paragraphs:
                            self.logger.warning(f"Word document has more than {self.max_word_paragraphs} paragraphs, truncating")
                            break

                # FIX: Extract structured tables using helper (TableEntityExtractor-compatible format)
                tables = self._extract_tables_word(doc)

                # Add table text to extracted_text for content search
                for table in doc.tables:
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            extracted_text.append('\t'.join(row_text))

                return {
                    'status': 'completed',
                    'method': 'word_native',
                    'text': '\n'.join(extracted_text),
                    'data': {'tables': tables},  # FIX: Add structured tables
                    'ocr_confidence': 1.0
                }
                
            finally:
                os.unlink(tmp_path)
                
        except Exception as e:
            return {
                'status': 'error',
                'method': 'word_failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_powerpoint(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process PowerPoint presentations"""
        try:
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                prs = Presentation(tmp_path)
                extracted_text = []
                slide_count = 0

                for slide_num, slide in enumerate(prs.slides):
                    # Enforce slide limit to prevent memory exhaustion
                    if slide_count >= self.max_ppt_slides:
                        self.logger.warning(f"PowerPoint has more than {self.max_ppt_slides} slides, truncating")
                        break

                    slide_count += 1
                    slide_text = [f"=== Slide {slide_num + 1} ==="]

                    # Extract text from shapes
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text.append(shape.text.strip())

                    # Extract notes
                    if slide.notes_slide and slide.notes_slide.notes_text_frame:
                        notes_text = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_text:
                            slide_text.append(f"[Notes: {notes_text}]")

                    extracted_text.extend(slide_text)

                # FIX: Extract structured tables using helper (TableEntityExtractor-compatible format)
                tables = self._extract_tables_powerpoint(prs)

                return {
                    'status': 'completed',
                    'method': 'powerpoint_native',
                    'text': '\n'.join(extracted_text),
                    'data': {'tables': tables},  # FIX: Add structured tables
                    'page_count': slide_count,
                    'ocr_confidence': 1.0
                }
                
            finally:
                os.unlink(tmp_path)
                
        except Exception as e:
            return {
                'status': 'error',
                'method': 'powerpoint_failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_image(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process images using OCR"""
        try:
            if not self.ocr_engine:
                return {
                    'status': 'error',
                    'method': 'ocr_unavailable',
                    'text': '',
                    'error': 'OCR engine not available'
                }
            
            # Save to temporary file
            with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as tmp_file:
                tmp_file.write(content)
                tmp_path = tmp_file.name
            
            try:
                ocr_result = self.ocr_engine.process_image(tmp_path)
                return {
                    'status': 'completed',
                    'method': 'image_ocr',
                    'text': ocr_result.get('text', ''),
                    'ocr_confidence': ocr_result.get('confidence', 0.0)
                }
                
            finally:
                os.unlink(tmp_path)
                
        except Exception as e:
            return {
                'status': 'error',
                'method': 'image_failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_text(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            text = content.decode('utf-8', errors='ignore')
            return {
                'status': 'completed',
                'method': 'text_native',
                'text': text,
                'ocr_confidence': 1.0
            }
        except Exception as e:
            return {
                'status': 'error',
                'method': 'text_failed',
                'text': '',
                'error': str(e)
            }
    
    def _process_csv(self, content: bytes, filename: str, storage_dir: Path) -> Dict[str, Any]:
        """Process CSV files"""
        try:
            import csv
            import io
            
            text_content = content.decode('utf-8', errors='ignore')
            csv_reader = csv.reader(io.StringIO(text_content))
            
            extracted_lines = []
            for row_num, row in enumerate(csv_reader):
                if row_num > 1000:  # Limit rows to prevent excessive processing
                    break
                extracted_lines.append('\t'.join(row))
            
            return {
                'status': 'completed',
                'method': 'csv_native',
                'text': '\n'.join(extracted_lines),
                'ocr_confidence': 1.0
            }
        except Exception as e:
            return {
                'status': 'error',
                'method': 'csv_failed',
                'text': '',
                'error': str(e)
            }
    
    def get_processing_stats(self) -> Dict[str, Any]:
        """Get processing statistics"""
        try:
            total_files = 0
            total_size = 0
            file_types = {}
            
            for year_dir in self.storage_path.iterdir():
                if year_dir.is_dir():
                    for month_dir in year_dir.iterdir():
                        if month_dir.is_dir():
                            for day_dir in month_dir.iterdir():
                                if day_dir.is_dir():
                                    for hash_dir in day_dir.iterdir():
                                        if hash_dir.is_dir():
                                            original_dir = hash_dir / 'original'
                                            if original_dir.exists():
                                                for file_path in original_dir.iterdir():
                                                    if file_path.is_file():
                                                        total_files += 1
                                                        total_size += file_path.stat().st_size
                                                        
                                                        # Count file types
                                                        suffix = file_path.suffix.lower()
                                                        file_types[suffix] = file_types.get(suffix, 0) + 1
            
            return {
                'total_files': total_files,
                'total_size_mb': total_size / (1024 * 1024),
                'file_types': file_types,
                'storage_path': str(self.storage_path)
            }
            
        except Exception as e:
            self.logger.error(f"Error getting processing stats: {e}")
            return {}
    
    def cleanup_old_files(self, days_old: int = 90):
        """Clean up old processed files"""
        try:
            import shutil
            from datetime import datetime, timedelta
            
            cutoff_date = datetime.now() - timedelta(days=days_old)
            cleaned_count = 0
            
            for year_dir in self.storage_path.iterdir():
                if year_dir.is_dir() and year_dir.name.isdigit():
                    year = int(year_dir.name)
                    if year < cutoff_date.year:
                        shutil.rmtree(year_dir)
                        cleaned_count += 1
                        self.logger.info(f"Cleaned up directory: {year_dir}")
            
            self.logger.info(f"Cleaned up {cleaned_count} old directories")
            
        except Exception as e:
            self.logger.error(f"Error cleaning up old files: {e}")